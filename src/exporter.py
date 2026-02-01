import pandas as pd
from pptx import Presentation

def export_excel(df, path):
    with pd.ExcelWriter(path, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False, sheet_name="Full Analysis")

        sector_summary = df.groupby("Sector")["Total Allocation (₹ Cr)"].sum()
        sector_summary.to_excel(writer, sheet_name="Sector Summary")

def export_ppt(report_text, path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Budget Allocation Optimality Report"
    slide.placeholders[1].text = report_text[:1200]
    prs.save(path)
